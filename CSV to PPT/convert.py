from pptx import Presentation
import csv
from pptx.util import Inches
from pptx.dml.color import RGBColor


prs = Presentation('empty1.pptx')
title_slide_layout = prs.slide_layouts[0]
title_only_slide_layout = prs.slide_layouts[3]

slid = prs.slides.add_slide(title_slide_layout)
title = slid.shapes.title
subtitle = slid.placeholders[1]

title.text = "IPL Auction 2021"
subtitle.text = "Oculus"



heading = []

with open('IPL.csv') as csv_file:
    csv_reader = csv.reader(csv_file, delimiter=',')
    line_count = 0
    for row in csv_reader:
        if line_count == 0:
            print(f'Column names are {", ".join(row)}')
            heading = row
            print(row)
            line_count += 1
        else:
            slide = prs.slides.add_slide(title_only_slide_layout)
            # for shape in slide.placeholders:
            #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
            # placeholder = slide.placeholders[13]
            shapes = slide.shapes
            placeholderpic = slide.placeholders[14]
            
            # print(row[0]+'.png')
            try:
                picture = placeholderpic.insert_picture('Pics/'+ row[0]+'.png')
                # fill = placeholderpic.fill
                # fill.gradient()
                # gradient_stops = fill.gradient_stops
                # gradient_stop = gradient_stops[0]
                # color = gradient_stop.color
                # color.rgb = RGB(0x3F, 0x2c, 0x36)

            except:
                print(row[0],"Not found")

            titletext = row[0]
            if row[3] == '1':
                titletext+=' (Wk)'
            if row[4] == '1':
                titletext+=' (UC)'
            shapes.title.text = titletext
            # graphic_frame = placeholder.insert_table(rows=13, cols=2)
            # table = graphic_frame.table
            # # rows = 13
            # # cols = 2
            # table.left = top = Inches(2.0)
            # table.width = Inches(6.0)
            # table.height = Inches(0.18)

            # # table = shapes.add_table(rows, cols, left, top, width, height).table
            # table.first_row = False
            # table.first_col = True
            # # set column widths
            # table.columns[0].width = Inches(3.0)
            # table.columns[1].width = Inches(3.0)

            
            # for i in range(1,14):
            #     table.cell(i-1, 0).text = heading[i].upper()
            #     table.cell(i-1, 1).text = row[i]
            
            #6 are internal stats
            plac = slide.placeholders[15]
            if row[8] == '0':
                plac.text = '-'
            else:   
                plac.text = str(row[8])
            
            plac = slide.placeholders[18]
            if row[9] == '0':
                plac.text = '-'
            else:   
                plac.text = str(row[9])
            
            
            plac = slide.placeholders[16]
            if row[10] == '0':
                plac.text = '-'
            else:   
                plac.text = str(row[10])
            
            
            plac = slide.placeholders[19]
            if row[11] == '0':
                plac.text = '-'
            else:   
                plac.text = str(row[11])
            
            
            plac = slide.placeholders[17]
            if row[12] == '0':
                plac.text = '-'
            else:   
                plac.text = str(row[12])
            
            
            plac = slide.placeholders[20]
            if row[13] == '0':
                plac.text = '-'
            else:   
                plac.text = str(row[13])
            
            #Overall
            plac = slide.placeholders[21]
            plac.text = str(row[7])

            #Pics
            plac = slide.placeholders[22]
            if row[1] == '0':
                plac.text = '-'
            else:   
                pic = plac.insert_picture('logo/'+ row[1]+'.png')
                # pic.width = Inches(1.0)
            
            plac = slide.placeholders[23]
            if row[6] == '1':
                pic = plac.insert_picture('plane.png')
            else:
                pic = plac.insert_picture('india.png')
             
            #starred
            plac = slide.placeholders[24]
            if row[5] == '1':
                pic = plac.insert_picture('star.png')
                # fill = placeholderpic.fill
                # fill.solid()
                # fore_color = fill.fore_color
                # fore_color.rgb = RGBColor(0x00, 0x00, 0x00)

            plac = slide.placeholders[25]
            if row[2] == '0':
                plac.text = '-'
            else:   
                plac.text = str(row[2])

            plac = slide.placeholders[26]
            if row[16] == '0':
                plac.text = '-'
            else:   
                plac.text = str(row[16])
            
                

            line_count += 1
    print(f'Processed {line_count} lines.')
prs.save('test.pptx')
print('Done!!')